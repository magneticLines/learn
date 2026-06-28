# -*- coding: utf-8 -*-
"""生成大作业演示 PPT（题目 8：工业设备故障诊断多传感器融合）。

12 页工业风格演示，含图片、流程图、卡片与 PowerPoint 进入动画。
运行：python build_ppt.py
输出：../大作业答辩_工业故障诊断融合.pptx
"""
import os
import json
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

import ppt_helpers as H
from ppt_helpers import (add_text, add_rect, add_shape_text, add_circle_num,
                         add_pic, add_pic_h, fill_bg, animate)

ROOT = os.path.dirname(os.path.abspath(__file__))
FIG = os.path.join(ROOT, "figures")
OUT = os.path.abspath(os.path.join(ROOT, "..", "大作业答辩_工业故障诊断融合.pptx"))
M = json.load(open(os.path.join(ROOT, "results", "metrics.json"), encoding="utf-8"))
ACC = M["accuracy"]
FUSED = ACC["三传感器融合"]
BEST = max(ACC["仅振动"], ACC["仅声学"], ACC["仅温度"])
GAIN = (FUSED - BEST) * 100

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]
W = 13.333


def slide():
    return prs.slides.add_slide(BLANK)


def title_block(s, text, color=H.INK):
    return add_text(s, 0.6, 0.45, 12.1, 0.95, text, size=30, color=color,
                    bold=True, font=H.HEAD_FONT)


def kicker(s, text, color=H.TEAL):
    return add_text(s, 0.62, 0.3, 8, 0.4, text, size=12.5, color=color,
                    bold=True, font=H.BODY_FONT)


# ============ S1 封面（深色）============
s = slide(); fill_bg(s, H.BG_DARK)
deco = add_rect(s, 9.7, -1.6, 5.2, 5.2, fill=H.PRIMARY, shape=MSO_SHAPE.OVAL)
deco2 = add_rect(s, 11.2, 4.4, 3.6, 3.6, fill=H.TEAL, shape=MSO_SHAPE.OVAL)
t1 = add_text(s, 0.9, 2.0, 9.8, 0.5, "《多传感器数据融合》课程大作业", size=18,
              color=H.RGBColor(0xCA, 0xDC, 0xFC), bold=True)
t2 = add_text(s, 0.9, 2.6, 11.0, 1.7, "工业设备故障诊断\n多传感器融合设计与实现", size=40,
              color=H.WHITE, bold=True, font=H.HEAD_FONT, line_spacing=1.05)
t3 = add_text(s, 0.95, 4.7, 10.5, 0.5, "基于 振动 + 声学 + 红外温度 的滚动轴承故障诊断",
              size=16, color=H.RGBColor(0x9F, 0xC0, 0xD4), italic=True)
t4 = add_text(s, 0.95, 6.5, 11.0, 0.5, "汇报人：________      指导教师：________      2026 年 6 月",
              size=13, color=H.RGBColor(0x8F, 0xA6, 0xB8))
animate(s, [(t1, "fade"), (t2, "fade"), (t3, "wipe_right"), (t4, "fade")])

# ============ S2 目录（浅色）============
s = slide(); fill_bg(s, H.LIGHT)
kicker(s, "CONTENTS")
title_block(s, "目录")
items = ["研究背景与意义", "传感器配置与思路", "系统总体框架", "数据生成与特征提取",
         "实验结果与分析", "结论与创新点"]
anim = []
for i, it in enumerate(items):
    col = i // 3; row = i % 3
    x = 0.9 + col * 6.2; y = 2.0 + row * 1.45
    card = add_rect(s, x, y, 5.6, 1.15, fill=H.CARD, shadow=True)
    cnum = add_circle_num(s, x + 0.28, y + 0.3, 0.55, i + 1, fill=H.ACCENT)
    txt = add_text(s, x + 1.05, y, 4.4, 1.15, it, size=18, color=H.INK, bold=True,
                   anchor=MSO_ANCHOR.MIDDLE)
    anim += [(card, "fade"), (cnum, "circle"), (txt, "wipe_right")]
animate(s, anim)

# ============ S3 研究背景与意义 ============
s = slide(); fill_bg(s, H.LIGHT)
kicker(s, "01  研究背景")
title_block(s, "为什么需要多传感器融合？")
bg = add_text(s, 0.7, 1.85, 5.7, 4.6,
              "工业旋转设备的滚动轴承长期运行在高速、重载、变工况环境下，"
              "是故障高发部位；故障若不能及时发现，轻则停机停产，重则引发安全事故。\n\n"
              "传统诊断多依赖单一振动传感器，信息维度有限、易受工况与噪声干扰，"
              "在复杂场景下难以兼顾准确率与鲁棒性。\n\n"
              "多传感器融合通过整合互补信息，突破单一传感器的局限。",
              size=15.5, color=H.INK, line_spacing=1.18)
cards = [("振动", "对故障类型敏感，\n但易受机械噪声干扰"),
         ("声学", "便于非接触监测，\n但难以分辨故障类型"),
         ("温度", "反映整体热状态，\n但存在响应滞后")]
anim = [(bg, "fade")]
for i, (h, d) in enumerate(cards):
    y = 1.95 + i * 1.5
    c = add_rect(s, 7.0, y, 5.6, 1.3, fill=H.CARD, shadow=True)
    dot = add_circle_num(s, 7.25, y + 0.32, 0.62, "", fill=H.TEAL)
    add_shape_text(dot, h, size=14, color=H.WHITE, bold=True)
    add_text(s, 8.15, y + 0.12, 4.3, 1.1, d, size=13.5, color=H.INK,
             anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.05)
    anim += [(c, "wipe_up")]
animate(s, anim)

# ============ S4 传感器配置与思路 ============
s = slide(); fill_bg(s, H.LIGHT)
kicker(s, "02  传感器配置")
title_block(s, "三传感器各有所长 · 信息互补")
sub = add_text(s, 0.7, 1.55, 12, 0.5,
               "题目 8：基于振动 + 声学 + 红外温度的滚动轴承故障诊断（正常 / 内圈 / 外圈 / 滚动体 四种工况）",
               size=14, color=H.MUTED)
cfg = [("振动传感器", "擅长：故障类型", "时域 RMS/峭度 + 频域\nBPFI/BPFO/BSF 特征频率", H.PRIMARY),
       ("声学麦克风", "擅长：有无故障", "频谱熵 / 高频能量占比\n/ 异响频带能量", H.TEAL),
       ("红外温度", "擅长：严重程度", "均值 / 最大温升梯度\n/ 热点数", H.ACCENT)]
anim = [(sub, "fade")]
for i, (h, role, feat, col) in enumerate(cfg):
    x = 0.7 + i * 4.15
    card = add_rect(s, x, 2.3, 3.85, 4.1, fill=H.CARD, shadow=True)
    icon = add_rect(s, x + 1.5, 2.65, 0.85, 0.85, fill=col, shape=MSO_SHAPE.OVAL)
    add_shape_text(icon, str(i + 1), size=22, color=H.WHITE, bold=True)
    add_text(s, x + 0.15, 3.7, 3.55, 0.5, h, size=19, color=H.INK, bold=True,
             align=PP_ALIGN.CENTER)
    add_text(s, x + 0.15, 4.25, 3.55, 0.45, role, size=15, color=col, bold=True,
             align=PP_ALIGN.CENTER)
    add_text(s, x + 0.2, 4.9, 3.45, 1.3, feat, size=13, color=H.MUTED,
             align=PP_ALIGN.CENTER, line_spacing=1.15)
    anim += [(card, "fade"), (icon, "circle")]
animate(s, anim)

# ============ S5 总体框架 ============
s = slide(); fill_bg(s, H.LIGHT)
kicker(s, "03  总体框架")
title_block(s, "系统总体流程")
steps = [("数据采集", "三传感器\n并行采集", H.PRIMARY),
         ("特征提取", "时域/频域\n/温度特征", H.TEAL),
         ("特征级融合", "18 维特征\n拼接", H.ACCENT),
         ("随机森林", "集成学习\n分类", H.TEAL),
         ("诊断结果", "四工况\n输出", H.PRIMARY)]
anim = []
bw, gap = 2.15, 0.35
x0 = (W - (len(steps) * bw + (len(steps) - 1) * gap)) / 2
for i, (h, d, col) in enumerate(steps):
    x = x0 + i * (bw + gap)
    box = add_rect(s, x, 3.0, bw, 1.9, fill=col, shadow=True)
    add_shape_text(box, h + "\n", size=17, color=H.WHITE, bold=True)
    add_text(s, x + 0.1, 4.0, bw - 0.2, 0.8, d, size=12, color=H.WHITE,
             align=PP_ALIGN.CENTER, line_spacing=1.0)
    anim.append((box, "wipe_right"))
    if i < len(steps) - 1:
        arr = add_rect(s, x + bw + 0.02, 3.72, gap - 0.04, 0.46, fill=H.MUTED,
                       shape=MSO_SHAPE.RIGHT_ARROW)
        anim.append((arr, "fade"))
note = add_text(s, 0.7, 5.7, 12, 0.6,
                "三类传感器并行采集 → 各自特征提取 → 拼接为统一特征向量 → 随机森林输出诊断结果",
                size=14, color=H.MUTED, align=PP_ALIGN.CENTER)
anim.append((note, "fade"))
animate(s, anim)

# ============ S6 数据生成 ============
s = slide(); fill_bg(s, H.LIGHT)
kicker(s, "04  数据生成")
title_block(s, "多传感器数据生成（机理模拟）")
img = add_pic(s, os.path.join(FIG, "fig1_vibration_waveform.png"), 0.6, 1.95, 7.85)
pts = add_text(s, 8.7, 2.1, 4.2, 4.7,
               "无真实传感器，按物理机理模拟生成：\n\n"
               "• 振动：转频谐波 + 故障特征频率冲击\n   + 包络谱线 + 噪声\n\n"
               "• 声学：声辐射 + 故障异响（各类型相近）\n   + 环境噪声\n\n"
               "• 温度：一阶温升稳态 + 故障热点\n\n"
               "4 工况 × 200 = 800 样本，\n三传感器可分性各有侧重且存在重叠。",
               size=14.5, color=H.INK, line_spacing=1.22)
animate(s, [(img, "fade"), (pts, "wipe_right")])

# ============ S7 特征体系 ============
s = slide(); fill_bg(s, H.LIGHT)
kicker(s, "04  特征提取")
title_block(s, "18 维特征体系")
groups = [("振动 (9维)", H.PRIMARY,
           "时域：RMS、峭度、峰值因子、标准差\n频域：BPFI/BPFO/BSF 频带能量、频谱重心"),
          ("声学 (4维)", H.TEAL,
           "RMS、频谱熵、\n高频能量占比、异响频带能量"),
          ("温度 (5维)", H.ACCENT,
           "均值、最大值、\n最大温升梯度、总温升、热点数")]
anim = []
for i, (h, col, body) in enumerate(groups):
    x = 0.7 + i * 4.15
    card = add_rect(s, x, 2.0, 3.85, 3.9, fill=H.CARD, shadow=True)
    hd = add_rect(s, x + 0.3, 2.3, 3.25, 0.7, fill=col)
    add_shape_text(hd, h, size=16, color=H.WHITE, bold=True)
    add_text(s, x + 0.35, 3.35, 3.2, 2.3, body, size=13.5, color=H.INK, line_spacing=1.25)
    anim += [(card, "fade"), (hd, "wipe_right")]
animate(s, anim)

# ============ S8 融合方法 ============
s = slide(); fill_bg(s, H.LIGHT)
kicker(s, "04  融合方法")
title_block(s, "特征级融合 + 随机森林分类")
anim = []
srcs = [("振动特征", H.PRIMARY), ("声学特征", H.TEAL), ("温度特征", H.ACCENT)]
for i, (h, col) in enumerate(srcs):
    b = add_rect(s, 0.8, 2.0 + i * 1.25, 2.6, 1.0, fill=col, shadow=True)
    add_shape_text(b, h, size=15, color=H.WHITE, bold=True)
    anim.append((b, "wipe_right"))
merge = add_rect(s, 4.3, 2.9, 2.7, 1.2, fill=H.INK, shadow=True)
add_shape_text(merge, "18 维\n融合特征", size=15, color=H.WHITE, bold=True)
rf = add_rect(s, 7.8, 2.9, 2.5, 1.2, fill=H.TEAL, shadow=True)
add_shape_text(rf, "随机森林\n200 棵树", size=15, color=H.WHITE, bold=True)
out = add_rect(s, 11.0, 2.9, 1.9, 1.2, fill=H.ACCENT, shadow=True)
add_shape_text(out, "四工况\n诊断", size=15, color=H.WHITE, bold=True)
a1 = add_rect(s, 3.45, 3.35, 0.8, 0.4, fill=H.MUTED, shape=MSO_SHAPE.RIGHT_ARROW)
a2 = add_rect(s, 7.05, 3.35, 0.7, 0.4, fill=H.MUTED, shape=MSO_SHAPE.RIGHT_ARROW)
a3 = add_rect(s, 10.35, 3.35, 0.6, 0.4, fill=H.MUTED, shape=MSO_SHAPE.RIGHT_ARROW)
cap = add_text(s, 0.8, 6.05, 11.7, 0.6,
               "特征级融合：保留各传感器细节信息，在信息保留与实现复杂度之间取得平衡。",
               size=14, color=H.MUTED, align=PP_ALIGN.CENTER)
anim += [(merge, "fade"), (a1, "fade"), (rf, "fade"), (a2, "fade"),
         (out, "fade"), (a3, "fade"), (cap, "fade")]
animate(s, anim)

# ============ S9 结果-准确率 ============
s = slide(); fill_bg(s, H.LIGHT)
kicker(s, "05  实验结果")
title_block(s, "融合显著优于任一单传感器")
img = add_pic(s, os.path.join(FIG, "fig4_accuracy_comparison.png"), 0.6, 1.75, 6.6)
big = add_text(s, 7.7, 1.9, 5.2, 1.5, f"{FUSED*100:.1f}%", size=72, color=H.ACCENT,
               bold=True, font=H.HEAD_FONT)
lbl = add_text(s, 7.75, 3.35, 5.2, 0.5, "三传感器融合准确率", size=17, color=H.INK, bold=True)
comp = add_text(s, 7.75, 4.1, 5.3, 2.0,
                f"仅振动：{ACC['仅振动']*100:.1f}%\n"
                f"仅声学：{ACC['仅声学']*100:.1f}%\n"
                f"仅温度：{ACC['仅温度']*100:.1f}%",
                size=16, color=H.MUTED, line_spacing=1.3)
gain = add_text(s, 7.75, 6.0, 5.3, 0.7, f"较最佳单传感器  ▲ +{GAIN:.1f} 个百分点",
                size=16, color=H.TEAL, bold=True)
animate(s, [(img, "fade"), (big, "circle"), (lbl, "fade"), (comp, "wipe_up"), (gain, "fade")])

# ============ S10 结果-混淆矩阵&特征重要性 ============
s = slide(); fill_bg(s, H.LIGHT)
kicker(s, "05  结果分析")
title_block(s, "混淆矩阵 & 特征重要性")
im1 = add_pic_h(s, os.path.join(FIG, "fig5_confusion_matrices.png"), 0.7, 1.75, 4.75)
im2 = add_pic(s, os.path.join(FIG, "fig6_feature_importance.png"), 6.7, 2.0, 6.0)
c1 = add_text(s, 0.7, 6.55, 5.5, 0.4, "融合后各类别对角线占比显著提高，误判明显减少",
              size=12.5, color=H.MUTED)
c2 = add_text(s, 6.7, 6.55, 6.0, 0.4, "三类传感器特征均位居重要性前列，互补贡献",
              size=12.5, color=H.MUTED)
animate(s, [(im1, "fade"), (c1, "fade"), (im2, "fade"), (c2, "fade")])

# ============ S11 结论与创新点（深色）============
s = slide(); fill_bg(s, H.BG_DARK)
add_text(s, 0.62, 0.3, 8, 0.4, "06  CONCLUSION", size=12.5, color=H.RGBColor(0x9F,0xC0,0xD4), bold=True)
title_block(s, "结论与创新点", color=H.WHITE)
inno = ["构建振动—声学—红外温度三类异构传感器统一的特征级融合诊断框架",
        "设计体现“类型—有无—严重度”互补性的 18 维特征体系",
        "从准确率、混淆矩阵、特征重要性三角度量化验证融合有效性"]
anim = []
for i, it in enumerate(inno):
    y = 1.9 + i * 1.0
    cnum = add_circle_num(s, 0.8, y, 0.6, i + 1, fill=H.ACCENT)
    txt = add_text(s, 1.7, y - 0.02, 8.0, 0.9, it, size=16,
                   color=H.RGBColor(0xE6,0xEE,0xF4), anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.05)
    anim += [(cnum, "circle"), (txt, "wipe_right")]
stat = add_rect(s, 10.0, 2.0, 2.7, 2.9, fill=H.PRIMARY, shadow=True)
add_shape_text(stat, f"{FUSED*100:.1f}%\n融合准确率\n\n▲ +{GAIN:.1f} pts", size=20,
               color=H.WHITE, bold=True, line_spacing=1.1)
anim.append((stat, "fade"))
animate(s, anim)

# ============ S12 谢谢（深色）============
s = slide(); fill_bg(s, H.BG_DARK)
d1 = add_rect(s, -1.5, 4.6, 4.6, 4.6, fill=H.PRIMARY, shape=MSO_SHAPE.OVAL)
d2 = add_rect(s, 10.6, -1.6, 4.4, 4.4, fill=H.TEAL, shape=MSO_SHAPE.OVAL)
ty = add_text(s, 0, 2.6, W, 1.3, "谢 谢 聆 听", size=54, color=H.WHITE, bold=True,
              font=H.HEAD_FONT, align=PP_ALIGN.CENTER)
tt = add_text(s, 0, 4.2, W, 0.6, "敬请各位老师批评指正", size=20,
              color=H.RGBColor(0x9F,0xC0,0xD4), align=PP_ALIGN.CENTER, italic=True)
animate(s, [(ty, "fade"), (tt, "fade")])

prs.save(OUT)
print("PPT 已生成:", OUT, "| 幻灯片数:", len(prs.slides._sldIdLst))
