# -*- coding: utf-8 -*-
"""PPT 构建辅助：颜色、文本、形状、图片，以及 PowerPoint 进入动画注入。"""
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from copy import deepcopy
import lxml.etree as etree

# ---- 配色（工业/工程 深蓝-teal 主题）----
BG_DARK = RGBColor(0x0E, 0x2A, 0x3B)
PRIMARY = RGBColor(0x06, 0x5A, 0x82)
TEAL = RGBColor(0x1C, 0x72, 0x93)
LIGHT = RGBColor(0xF4, 0xF7, 0xFA)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
INK = RGBColor(0x1B, 0x2A, 0x36)
MUTED = RGBColor(0x6B, 0x7C, 0x8A)
ACCENT = RGBColor(0xE1, 0x65, 0x5B)     # 珊瑚红，用于“融合”强调，与图表配色一致
CARD = RGBColor(0xFF, 0xFF, 0xFF)
SAND = RGBColor(0xE9, 0xEF, 0xF4)

HEAD_FONT = "Cambria"      # 安全字体（标题，带衬线个性）
BODY_FONT = "Calibri"      # 安全字体（正文）


def fill_bg(slide, color):
    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = color


def _set_font(run, size, color, bold, font, italic=False):
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.italic = italic
    run.font.name = font
    # 中文 eastAsia 字体（用微软雅黑，确保中文显示）
    rpr = run._r.get_or_add_rPr()
    ea = rpr.find(qn("a:ea"))
    if ea is None:
        ea = etree.SubElement(rpr, qn("a:ea"))
    ea.set("typeface", "微软雅黑")


def add_text(slide, l, t, w, h, text, size=16, color=INK, bold=False,
             font=BODY_FONT, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
             italic=False, line_spacing=1.0):
    box = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = 0; tf.margin_right = 0
    tf.margin_top = 0; tf.margin_bottom = 0
    lines = text.split("\n")
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        run = p.add_run(); run.text = ln
        _set_font(run, size, color, bold, font, italic)
    return box


def add_rect(slide, l, t, w, h, fill=CARD, line=None, shape=MSO_SHAPE.ROUNDED_RECTANGLE,
             line_w=1.0, shadow=False):
    sp = slide.shapes.add_shape(shape, Inches(l), Inches(t), Inches(w), Inches(h))
    sp.fill.solid(); sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line; sp.line.width = Pt(line_w)
    sp.shadow.inherit = False
    if shadow:
        _soft_shadow(sp)
    # 去掉圆角默认过大：调小圆角半径
    if shape == MSO_SHAPE.ROUNDED_RECTANGLE:
        try:
            sp.adjustments[0] = 0.06
        except Exception:
            pass
    return sp


def _soft_shadow(sp):
    spPr = sp._element.spPr
    # 复用 shadow.inherit=False 生成的空 effectLst，避免出现两个同级 effectLst（PowerPoint 会拒绝）
    effLst = spPr.find(qn("a:effectLst"))
    if effLst is None:
        effLst = etree.SubElement(spPr, qn("a:effectLst"))
    else:
        for ch in list(effLst):
            effLst.remove(ch)
    outer = etree.SubElement(effLst, qn("a:outerShdw"))
    outer.set("blurRad", "90000"); outer.set("dist", "38100")
    outer.set("dir", "5400000"); outer.set("rotWithShape", "0")
    clr = etree.SubElement(outer, qn("a:srgbClr")); clr.set("val", "1B2A36")
    alpha = etree.SubElement(clr, qn("a:alpha")); alpha.set("val", "24000")


def add_shape_text(sp, text, size=16, color=WHITE, bold=True, font=BODY_FONT,
                   align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.0):
    tf = sp.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = Inches(0.08); tf.margin_right = Inches(0.08)
    tf.margin_top = Inches(0.04); tf.margin_bottom = Inches(0.04)
    for i, ln in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align; p.line_spacing = line_spacing
        run = p.add_run(); run.text = ln
        _set_font(run, size, color, bold, font)
    return sp


def add_circle_num(slide, l, t, d, num, fill=ACCENT, txt=WHITE, size=18):
    c = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(l), Inches(t), Inches(d), Inches(d))
    c.fill.solid(); c.fill.fore_color.rgb = fill
    c.line.fill.background(); c.shadow.inherit = False
    add_shape_text(c, str(num), size=size, color=txt, bold=True)
    return c


def add_pic(slide, path, l, t, w):
    return slide.shapes.add_picture(path, Inches(l), Inches(t), width=Inches(w))


def add_pic_h(slide, path, l, t, h):
    return slide.shapes.add_picture(path, Inches(l), Inches(t), height=Inches(h))


# ----------------- 动画注入 -----------------
# 安全的 animEffect 进入滤镜（无需运动路径，兼容性好）
FILTERS = {
    "fade": "fade",
    "wipe_up": "wipe(up)",
    "wipe_right": "wipe(right)",
    "blinds": "blinds(horizontal)",
    "dissolve": "dissolve",
    "circle": "circle",
    "wedge": "wedge",
}


def _par_for_shape(spid, filt, ids):
    """生成单个形状的 clickEffect 进入动画 par 节点 XML 字符串。"""
    a, b, c, d, e = ids
    flt = FILTERS.get(filt, "fade")
    return f'''<p:par xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main">
  <p:cTn id="{a}" fill="hold">
    <p:stCondLst><p:cond delay="indefinite"/></p:stCondLst>
    <p:childTnLst>
      <p:par>
        <p:cTn id="{b}" fill="hold">
          <p:stCondLst><p:cond delay="0"/></p:stCondLst>
          <p:childTnLst>
            <p:par>
              <p:cTn id="{c}" presetID="10" presetClass="entr" presetSubtype="0" fill="hold" grpId="0" nodeType="clickEffect">
                <p:stCondLst><p:cond delay="0"/></p:stCondLst>
                <p:childTnLst>
                  <p:set>
                    <p:cBhvr>
                      <p:cTn id="{d}" dur="1" fill="hold"><p:stCondLst><p:cond delay="0"/></p:stCondLst></p:cTn>
                      <p:tgtEl><p:spTgt spid="{spid}"/></p:tgtEl>
                      <p:attrNameLst><p:attrName>style.visibility</p:attrName></p:attrNameLst>
                    </p:cBhvr>
                    <p:to><p:strVal val="visible"/></p:to>
                  </p:set>
                  <p:animEffect transition="in" filter="{flt}">
                    <p:cBhvr>
                      <p:cTn id="{e}" dur="500"/>
                      <p:tgtEl><p:spTgt spid="{spid}"/></p:tgtEl>
                    </p:cBhvr>
                  </p:animEffect>
                </p:childTnLst>
              </p:cTn>
            </p:par>
          </p:childTnLst>
        </p:cTn>
      </p:par>
    </p:childTnLst>
  </p:cTn>
</p:par>'''


def animate(slide, items):
    """为 slide 注入点击逐步进入动画。

    items: list[(shape, filter_name)]，按出现顺序排列。
    """
    if not items:
        return
    pars = []
    nid = 3
    for shape, filt in items:
        ids = (nid, nid + 1, nid + 2, nid + 3, nid + 4)
        nid += 5
        pars.append(_par_for_shape(shape.shape_id, filt, ids))
    timing_xml = f'''<p:timing xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main">
  <p:tnLst>
    <p:par>
      <p:cTn id="1" dur="indefinite" restart="never" nodeType="tmRoot">
        <p:childTnLst>
          <p:seq concurrent="1" nextAc="seek">
            <p:cTn id="2" dur="indefinite" nodeType="mainSeq">
              <p:childTnLst>
                {''.join(pars)}
              </p:childTnLst>
            </p:cTn>
            <p:prevCondLst><p:cond evt="onPrev" delay="0"><p:tgtEl><p:sldTgt/></p:tgtEl></p:cond></p:prevCondLst>
            <p:nextCondLst><p:cond evt="onNext" delay="0"><p:tgtEl><p:sldTgt/></p:tgtEl></p:cond></p:nextCondLst>
          </p:seq>
        </p:childTnLst>
      </p:cTn>
    </p:par>
  </p:tnLst>
</p:timing>'''
    timing = etree.fromstring(timing_xml)
    sld = slide._element            # <p:sld>
    # timing 必须是 p:sld 的最后一个子元素
    existing = sld.find(qn("p:timing"))
    if existing is not None:
        sld.remove(existing)
    sld.append(timing)
